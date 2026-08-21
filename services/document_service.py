"""
services/document_service.py
Extract text from PDF, DOCX, PPTX, TXT, and CSV files.
"""

import io
import pandas as pd


def extract_text(uploaded_file) -> dict:
    """
    Accepts a Streamlit UploadedFile and returns extracted text.

    Returns:
        {
            "filename": str,
            "file_type": str,
            "text": str,
            "word_count": int,
            "char_count": int,
        }
    """
    filename: str = uploaded_file.name
    ext = filename.rsplit(".", 1)[-1].lower()

    if ext == "pdf":
        text = _extract_pdf(uploaded_file)
    elif ext == "docx":
        text = _extract_docx(uploaded_file)
    elif ext == "pptx":
        text = _extract_pptx(uploaded_file)
    elif ext == "txt":
        text = uploaded_file.read().decode("utf-8", errors="ignore")
    elif ext == "csv":
        text = _extract_csv(uploaded_file)
    else:
        raise ValueError(f"❌ Unsupported file type: .{ext}")

    text = text.strip()
    if len(text) < 50:
        raise ValueError(f"⚠️ Could not extract meaningful text from {filename}.")

    return {
        "filename": filename,
        "file_type": ext.upper(),
        "text": text,
        "word_count": len(text.split()),
        "char_count": len(text),
    }


def _extract_pdf(uploaded_file) -> str:
    """Try pdfplumber first, fall back to PyPDF2."""
    data = uploaded_file.read()

    # pdfplumber
    try:
        import pdfplumber
        with pdfplumber.open(io.BytesIO(data)) as pdf:
            pages = [page.extract_text() or "" for page in pdf.pages]
        text = "\n".join(pages)
        if text.strip():
            return text
    except Exception:
        pass

    # PyPDF2 fallback
    try:
        import PyPDF2
        reader = PyPDF2.PdfReader(io.BytesIO(data))
        pages = [page.extract_text() or "" for page in reader.pages]
        return "\n".join(pages)
    except Exception as exc:
        raise ValueError(f"❌ Failed to read PDF: {exc}")


def _extract_docx(uploaded_file) -> str:
    try:
        from docx import Document
        doc = Document(uploaded_file)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        return "\n".join(paragraphs)
    except Exception as exc:
        raise ValueError(f"❌ Failed to read DOCX: {exc}")


def _extract_pptx(uploaded_file) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(uploaded_file)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
        return "\n".join(texts)
    except Exception as exc:
        raise ValueError(f"❌ Failed to read PPTX: {exc}")


def _extract_csv(uploaded_file) -> str:
    try:
        df = pd.read_csv(uploaded_file)
        return df.to_string(index=False)
    except Exception as exc:
        raise ValueError(f"❌ Failed to read CSV: {exc}")
